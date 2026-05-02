"""PPTX extraction: text via markitdown, images via python-pptx."""

from markitdown import MarkItDown
from pptx import Presentation
from pptx.util import Emu
from pathlib import Path


def extract_pptx(filepath, output_dir):
    output = Path(output_dir)
    (output / 'assets').mkdir(parents=True, exist_ok=True)

    # Text via markitdown
    md = MarkItDown()
    result = md.convert(filepath)
    content = result.text_content

    # Images from slide relationships
    prs = Presentation(filepath)
    img_num = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                img_num += 1
                image = shape.image
                ext = image.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                (output / 'assets' / f'figure-{img_num:02d}.{ext}').write_bytes(image.blob)

    # Multi-column layout reconstruction (if needed)
    # Group text shapes by slide, sort by shape.left position
    for slide in prs.slides:
        text_shapes = [
            (s.left, s.top, s.text_frame.text)
            for s in slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        text_shapes.sort(key=lambda t: (t[0] // Emu(3000000), t[1]))

    (output / 'content.md').write_text(content)
    return content
